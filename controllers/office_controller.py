import os
import io
import zipfile
import tempfile
import subprocess
import pandas as pd
from pdf2docx import Converter
import pdfplumber
from PIL import Image
import fitz
from pptx import Presentation
from pptx.util import Inches

class OfficeController:

    @staticmethod
    def pdf_to_pptx(file):
        doc = fitz.open(stream=file.read(), filetype="pdf")
        prs = Presentation()
        prs.slide_width = Inches(8.5)
        prs.slide_height = Inches(11)
        for page in doc:
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
            img_data = io.BytesIO(pix.tobytes("png"))
            slide.shapes.add_picture(img_data, 0, 0, width=prs.slide_width, height=prs.slide_height)
        out = io.BytesIO()
        prs.save(out)
        doc.close()
        out.seek(0)
        return out, "application/vnd.openxmlformats-officedocument.presentationml.presentation", "convertido.pptx"

    @staticmethod
    def pdf_to_html(file):
        doc = fitz.open(stream=file.read(), filetype="pdf")
        html_out = "<html><body>"
        for page in doc:
            html_out += page.get_text("html")
        html_out += "</body></html>"
        doc.close()
        out = io.BytesIO(html_out.encode('utf-8'))
        return out, "text/html", "convertido.html"

    @staticmethod
    def pdf_to_txt(file):
        doc = fitz.open(stream=file.read(), filetype="pdf")
        text_out = ""
        for page in doc:
            text_out += page.get_text()
        doc.close()
        out = io.BytesIO(text_out.encode('utf-8'))
        return out, "text/plain", "convertido.txt"

    @staticmethod
    def pdf_to_image(file, ext):
        doc = fitz.open(stream=file.read(), filetype="pdf")
        zip_buffer = io.BytesIO()
        fmt = "jpeg" if ext.upper() == "JPG" else "png"
        with zipfile.ZipFile(zip_buffer, "w") as zf:
            for i in range(len(doc)):
                page = doc.load_page(i)
                pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                zf.writestr(f"pagina_{i + 1}.{ext.lower()}", pix.tobytes(fmt))
        doc.close()
        zip_buffer.seek(0)
        return zip_buffer, "application/zip", f"imagenes.{ext.lower()}.zip"

    @staticmethod
    def pdf_to_word(file):
        with tempfile.TemporaryDirectory() as tmp_dir:
            pdf_p = os.path.join(tmp_dir, "in.pdf")
            docx_p = os.path.join(tmp_dir, "out.docx")
            with open(pdf_p, "wb") as f:
                f.write(file.read())
            
            cv = Converter(pdf_p)
            cv.convert(docx_p)
            cv.close()
            
            with open(docx_p, "rb") as f:
                out = io.BytesIO(f.read())
            return out, "application/vnd.openxmlformats-officedocument.wordprocessingml.document", "convertido.docx"

    @staticmethod
    def pdf_to_excel(file):
        tables = []
        with pdfplumber.open(io.BytesIO(file.read())) as pdf:
            for page in pdf.pages:
                tbl = page.extract_table()
                if tbl and len(tbl) > 1:
                    headers = tbl[0]
                    headers = [h if h else f"Col_{i}" for i, h in enumerate(headers)]
                    tables.append(pd.DataFrame(tbl[1:], columns=headers))
        if tables:
            out = io.BytesIO()
            pd.concat(tables, ignore_index=True).to_excel(out, index=False, engine='openpyxl')
            out.seek(0)
            return out, "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", "tablas.xlsx"
        else:
            raise ValueError("No se encontraron tablas estructuradas en el PDF.")

    @staticmethod
    def multiple_img_to_pdf(files):
        imgs = [Image.open(io.BytesIO(f.read())).convert('RGB') for f in files]
        out = io.BytesIO()
        imgs[0].save(out, format="PDF", save_all=True, append_images=imgs[1:])
        out.seek(0)
        return out

    @staticmethod
    def office_to_pdf(file, original_filename):
        with tempfile.TemporaryDirectory() as tmp_dir:
            in_p = os.path.join(tmp_dir, original_filename)
            with open(in_p, "wb") as f:
                f.write(file.read())
                
            comando = ["soffice", "--headless", "--convert-to", "pdf", in_p, "--outdir", tmp_dir]
            try:
                subprocess.run(comando, check=True, capture_output=True, timeout=60)
            except FileNotFoundError:
                raise Exception("El motor de LibreOffice no se encontró en el contenedor.")
                
            base_name = os.path.splitext(original_filename)[0]
            pdf_p = os.path.join(tmp_dir, f"{base_name}.pdf")
            
            with open(pdf_p, "rb") as f:
                out = io.BytesIO(f.read())
            return out
